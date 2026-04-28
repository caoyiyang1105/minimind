from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "实验3：大模型微调汇报.pptx"

FONT = "Microsoft YaHei"
TITLE = RGBColor(20, 32, 51)
TEXT = RGBColor(50, 58, 70)
MUTED = RGBColor(105, 116, 132)
BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(235, 243, 255)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
LINE = RGBColor(223, 230, 240)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)


def text_box(slide, x, y, w, h, text, size=22, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    text_box(slide, 0.55, 0.35, 8.8, 0.45, title, size=28, bold=True, color=TITLE)
    if subtitle:
        text_box(slide, 0.58, 0.83, 8.5, 0.28, subtitle, size=10.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.18), Inches(0.62), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def bullet_list(slide, x, y, w, h, items, size=17, color=TEXT, spacing=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(5 + spacing)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def pill(slide, x, y, w, h, text, fill=LIGHT_BLUE, color=BLUE, size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


def card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text_box(slide, x + 0.18, y + 0.18, w - 0.3, 0.28, title, size=14.5, bold=True, color=TITLE)
    text_box(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.65, body, size=11.5, color=TEXT)


def table(slide, x, y, w, h, data, widths=None, font_size=10.5):
    rows, cols = len(data), len(data[0])
    t = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        for i, col_w in enumerate(widths):
            t.columns[i].width = Inches(col_w)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = RGBColor(241, 245, 249) if r == 0 else RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = TITLE if r == 0 else TEXT
    return t


def add_image_fit(slide, path, x, y, w, h):
    path = ROOT / path
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        ww = w
        hh = w / img_ratio
    else:
        hh = h
        ww = h * img_ratio
    xx = x + (w - ww) / 2
    yy = y + (h - hh) / 2
    slide.shapes.add_picture(str(path), Inches(xx), Inches(yy), width=Inches(ww), height=Inches(hh))


def footer(slide, n):
    text_box(slide, 0.55, 7.0, 7.2, 0.2, "MiniMind 复现与《2025研究生手册》微调", size=8.5, color=MUTED)
    text_box(slide, 12.0, 7.0, 0.6, 0.2, str(n), size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    text_box(s, 0.72, 0.9, 11.5, 0.75, "实验3：大模型微调", size=40, bold=True, color=TITLE)
    text_box(s, 0.78, 1.85, 10.5, 0.45, "MiniMind mini 复现 + 官方模型对比 + 研究生手册领域微调", size=20, color=TEXT)
    pill(s, 0.8, 2.65, 1.35, 0.42, "Pretrain")
    pill(s, 2.32, 2.65, 0.95, 0.42, "SFT")
    pill(s, 3.45, 2.65, 1.38, 0.42, "Model Eval")
    pill(s, 5.0, 2.65, 1.55, 0.42, "HUST QA SFT")
    card(s, 0.8, 4.1, 3.5, 1.2, "实验要求", "复现 MiniMind 训练过程，给出 loss 曲线，并对比自训模型与官方模型效果。")
    card(s, 4.75, 4.1, 3.5, 1.2, "加分项", "利用《华中科技大学研究生手册》继续微调，并测试微调后的生成效果。", accent=GREEN)
    card(s, 8.7, 4.1, 3.2, 1.2, "汇报节奏", "8 页内容，约 5 分钟；重点看流程、曲线、效果和结论。", accent=ORANGE)
    footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实验路线", "从 mini 数据复现基础模型，再进行手册领域增强")
    steps = [
        ("1", "数据准备", "下载 mini 预训练/SFT 数据与官方 minimind-3"),
        ("2", "复现训练", "4 卡 DDP 完成 Pretrain 与 Full SFT"),
        ("3", "结果分析", "解析日志，绘制 loss 曲线，对比官方/自训生成效果"),
        ("4", "领域微调", "抽取《2025研究生手册》，构造 QA 数据继续 SFT"),
    ]
    x = 0.75
    for i, (num, title, body) in enumerate(steps):
        card(s, x + i * 3.1, 2.0, 2.55, 2.55, f"{num}. {title}", body, accent=BLUE if i < 3 else GREEN)
    bullet_list(s, 0.9, 5.25, 11.4, 0.9, [
        "主线：pretrain_t2t_mini → sft_t2t_mini → full_sft_mini_768.pth",
        "加分项：研究生手册 PDF → 文本抽取 → 手册 QA SFT → full_sft_hust_qa_768.pth",
    ], size=16)
    footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实验环境与数据", "硬件、数据规模和模型权重")
    table(s, 0.75, 1.65, 5.6, 2.25, [
        ["项目", "配置"],
        ["GPU", "4 × NVIDIA RTX 3090"],
        ["框架", "PyTorch 2.5.1+cu118"],
        ["环境", "Conda light, Python 3.10"],
        ["模型", "64M Dense minimind-3 Zero"],
    ], widths=[1.35, 4.15], font_size=12)
    table(s, 6.8, 1.65, 5.65, 2.25, [
        ["数据/模型", "规模"],
        ["pretrain_t2t_mini", "1,270,238 条 / 1.2GB"],
        ["sft_t2t_mini", "905,718 条 / 1.7GB"],
        ["官方 minimind-3", "122MB safetensors"],
        ["手册 QA SFT", "640 条"],
    ], widths=[2.7, 2.85], font_size=12)
    card(s, 0.75, 4.55, 5.6, 1.15, "数据来源", "MiniMind 官方 ModelScope 数据集；官方模型 minimind-3；华中科技大学研究生院《2025研究生手册》PDF。")
    card(s, 6.8, 4.55, 5.65, 1.15, "环境说明", "原计划 cu121，实际使用本机稳定可用的 cu118 环境；CUDA 可用并识别 4 张 3090。", accent=ORANGE)
    footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "训练设置", "复现 MiniMind mini 推荐训练流程")
    table(s, 0.75, 1.55, 11.9, 2.25, [
        ["阶段", "初始化", "Epoch", "Batch / Accum", "Max Len", "输出"],
        ["Pretrain", "none", "2", "8 / 8", "768", "pretrain_mini_768.pth"],
        ["Full SFT", "pretrain_mini", "2", "4 / 1", "768", "full_sft_mini_768.pth"],
        ["HUST QA SFT", "full_sft_mini", "8", "8 / 1", "768", "full_sft_hust_qa_768.pth"],
    ], widths=[1.55, 2.0, 1.0, 1.7, 1.2, 4.0], font_size=11.5)
    bullet_list(s, 1.0, 4.35, 11.4, 1.4, [
        "损失函数：自回归交叉熵；SFT 阶段只对 assistant 回复部分计算 label。",
        "日志字段：loss、logits_loss、aux_loss、lr；通过 plot_loss.py 输出 CSV 与 PNG。",
        "评估参数：固定 seed；官方/自训对比使用 temperature=0.7, top_p=0.9, max_new_tokens=512。",
    ], size=16)
    footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "训练损失曲线", "预训练和 SFT 均正常收敛")
    add_image_fit(s, "experiments/results/loss_curve.png", 0.65, 1.35, 7.0, 4.8)
    table(s, 8.05, 1.62, 4.55, 1.8, [
        ["阶段", "初始", "最终", "最低"],
        ["Pretrain", "7.3648", "1.6775", "1.3206"],
        ["SFT", "2.8538", "1.5734", "0.3544"],
    ], widths=[1.2, 1.1, 1.1, 1.1], font_size=12)
    card(s, 8.05, 4.0, 4.55, 1.45, "观察", "预训练从随机初始化快速下降；SFT 初始 loss 更低，说明模型已具备基础语言能力，之后学习指令跟随。")
    card(s, 8.05, 5.75, 4.55, 0.72, "结论", "mini 复现流程训练正常，权重和日志均完整保存。", accent=GREEN)
    footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "官方模型 vs 自训模型", "同题生成效果对比")
    table(s, 0.75, 1.5, 11.9, 2.55, [
        ["测试项", "官方 minimind-3", "自训 full_sft_mini"],
        ["自我介绍", "简短稳定，身份明确", "能说明 MiniMind 身份，但更冗长"],
        ["常识问答", "回答较集中，停止自然", "能答出要点，但有重复和概念混用"],
        ["机器学习", "结构较清晰", "内容更长，容易截断"],
        ["代码/计算", "基本可用", "基本可用，稳定性略弱"],
    ], widths=[1.7, 4.8, 5.4], font_size=10.8)
    card(s, 0.95, 4.55, 5.35, 1.2, "官方模型优势", "训练数据和流程更充分，回答更短、更稳，重复更少。")
    card(s, 6.95, 4.55, 5.35, 1.2, "自训模型结论", "完成了 mini 路线复现，具备基本问答能力；但仍弱于官方模型。", accent=ORANGE)
    text_box(s, 0.95, 6.25, 11.6, 0.3, "详细结果见 experiments/results/model_comparison.md", size=10.5, color=MUTED)
    footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "加分项：研究生手册微调", "用《2025研究生手册》增强领域问答")
    add_image_fit(s, "experiments/results/hust_qa_sft_loss_curve.png", 0.65, 1.35, 5.75, 3.6)
    table(s, 6.75, 1.45, 5.7, 1.6, [
        ["微调方式", "初始 loss", "最终 loss"],
        ["LoRA 条文", "3.0951", "2.5231"],
        ["全参条文 SFT", "2.3655", "0.3827"],
        ["最终 QA SFT", "0.6802", "0.0070"],
    ], widths=[2.3, 1.7, 1.7], font_size=11.5)
    card(s, 6.75, 3.65, 5.7, 1.05, "为什么用 QA SFT", "小模型容量有限，自动条文样本噪声较多；高质量 QA 更直接地告诉模型“怎么问、怎么答”。", accent=GREEN)
    table(s, 0.75, 5.35, 11.7, 1.05, [
        ["问题", "微调前", "微调后"],
        ["研究生请假规定", "幻觉出年龄限制/健康安全等", "准确回答病假证明、导师签署、院系审批、销假续假"],
        ["学籍管理适用对象", "泛化为教育部规定", "取得入学资格和具有学籍的各类研究生"],
    ], widths=[2.3, 4.4, 5.0], font_size=10.2)
    footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "结论与交付物", "5 分钟汇报总结")
    card(s, 0.75, 1.45, 3.65, 1.2, "训练复现", "Pretrain loss 7.3648 → 1.6775；SFT loss 2.8538 → 1.5734。", accent=BLUE)
    card(s, 4.85, 1.45, 3.65, 1.2, "模型对比", "自训模型具备基本能力，但官方模型生成更简洁稳定。", accent=ORANGE)
    card(s, 8.95, 1.45, 3.65, 1.2, "手册微调", "QA SFT 后可稳定回答指定手册问题，事实一致性明显提升。", accent=GREEN)
    bullet_list(s, 0.95, 3.25, 11.6, 1.4, [
        "主要权重：pretrain_mini_768.pth、full_sft_mini_768.pth、full_sft_hust_qa_768.pth",
        "主要报告：model_comparison.md、hust_qa_sft_comparison.md、实验3：大模型微调实验报告.pdf",
        "局限：小模型对长文档泛化有限；若要覆盖任意手册细节，建议结合 RAG 或扩充高质量 QA。",
    ], size=16)
    text_box(s, 0.95, 5.85, 11.5, 0.55, "谢谢！", size=28, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    footer(s, 8)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
